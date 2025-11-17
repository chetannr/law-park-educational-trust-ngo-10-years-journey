#!/usr/bin/env python3
"""
Extract images and text from PowerPoint presentation.
Creates a structured output for website development.
"""

import os
import json
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import hashlib

def sanitize_filename(text, max_length=50):
    """Create a safe filename from text."""
    if not text:
        return "untitled"
    # Remove special characters, keep only alphanumeric and spaces
    cleaned = "".join(c if c.isalnum() or c in (' ', '-', '_') else '' for c in text)
    # Replace spaces with underscores and limit length
    cleaned = cleaned.replace(' ', '_').strip('_')
    return cleaned[:max_length] or "untitled"

def extract_text_from_shape(shape):
    """Extract text from a shape."""
    if hasattr(shape, "text"):
        return shape.text.strip()
    return ""

def extract_images_from_shape(shape, slide_num, shape_index, output_dir):
    """Extract images from a shape and save them."""
    images = []
    
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        image = shape.image
        image_bytes = image.blob
        image_ext = image.ext
        
        # Generate a meaningful filename
        # Try to use alt text or shape name if available
        image_name = f"slide_{slide_num:02d}_image_{shape_index:02d}"
        if hasattr(shape, 'name') and shape.name:
            name_part = sanitize_filename(shape.name, 20)
            if name_part != "Picture":
                image_name = f"slide_{slide_num:02d}_{name_part}_{shape_index:02d}"
        
        filename = f"{image_name}.{image_ext}"
        filepath = output_dir / filename
        
        # Save image
        with open(filepath, 'wb') as f:
            f.write(image_bytes)
        
        images.append({
            "filename": filename,
            "path": str(filepath.relative_to(output_dir.parent)),
            "shape_index": shape_index
        })
    
    # Check for images in grouped shapes
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for i, sub_shape in enumerate(shape.shapes):
            sub_images = extract_images_from_shape(sub_shape, slide_num, f"{shape_index}_{i}", output_dir)
            images.extend(sub_images)
    
    return images

def extract_slide_content(pptx_path, output_base_dir):
    """Extract all content from PowerPoint presentation."""
    prs = Presentation(pptx_path)
    
    # Create output directories
    images_dir = output_base_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)
    
    slides_data = []
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        print(f"Processing slide {slide_num}...")
        
        slide_data = {
            "slide_number": slide_num,
            "text_content": [],
            "images": [],
            "all_text": ""
        }
        
        # Extract text from all shapes
        text_parts = []
        image_index = 0
        
        for shape_index, shape in enumerate(slide.shapes):
            # Extract text
            text = extract_text_from_shape(shape)
            if text:
                text_parts.append(text)
                slide_data["text_content"].append({
                    "shape_index": shape_index,
                    "text": text
                })
            
            # Extract images
            images = extract_images_from_shape(shape, slide_num, image_index, images_dir)
            if images:
                slide_data["images"].extend(images)
                image_index += len(images)
        
        # Combine all text
        slide_data["all_text"] = "\n\n".join(text_parts)
        
        # Generate slide title from first text element or use default
        if text_parts:
            first_text = text_parts[0].split('\n')[0].strip()
            slide_data["title"] = sanitize_filename(first_text[:50]) or f"slide_{slide_num}"
        else:
            slide_data["title"] = f"slide_{slide_num}"
        
        slides_data.append(slide_data)
        print(f"  - Found {len(slide_data['text_content'])} text elements and {len(slide_data['images'])} images")
    
    return slides_data

def main():
    """Main extraction function."""
    script_dir = Path(__file__).parent
    pptx_path = script_dir / "LAW PARK.pptx"
    output_dir = script_dir / "extracted_content"
    
    if not pptx_path.exists():
        print(f"Error: PowerPoint file not found at {pptx_path}")
        return
    
    print(f"Extracting content from: {pptx_path}")
    print(f"Output directory: {output_dir}")
    
    # Extract content
    slides_data = extract_slide_content(pptx_path, output_dir)
    
    # Save JSON mapping
    json_path = output_dir / "slides_data.json"
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(slides_data, f, indent=2, ensure_ascii=False)
    
    print(f"\nExtraction complete!")
    print(f"  - Total slides: {len(slides_data)}")
    print(f"  - Images saved to: {output_dir / 'images'}")
    print(f"  - Data saved to: {json_path}")
    
    # Print summary
    total_images = sum(len(slide['images']) for slide in slides_data)
    print(f"  - Total images extracted: {total_images}")

if __name__ == "__main__":
    main()

